"""
Builds the presenter slide-deck template (.potx): title slide, sponsor
"thank you" slide, and session-evaluation slide.

Driven entirely by a JSON manifest (see Generate-SlideTemplate.ps1, which
fetches event/sponsor data and calls this script). Not meant to be run by
hand except for testing -- use the .ps1 wrapper.
"""
import json
import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from PIL import Image, ImageChops

EMU_PER_IN = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

FONT = "Berlin Sans FB"


def emu(inches):
    return Emu(int(round(inches * EMU_PER_IN)))


def set_background_watermark(slide, image_path, alpha_pct=12):
    """Sets a slide's background to a picture, faded to alpha_pct opacity,
    stretched/cropped the way the previous year's template did (slightly
    cropped top and bottom so the logo reads large and centered)."""
    blip_fill = (
        f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f"<p:bgPr>"
        f'<a:blipFill dpi="0" rotWithShape="1">'
        f'<a:blip r:embed="{{RID}}"><a:alphaModFix amt="{alpha_pct * 1000}"/></a:blip>'
        f"<a:srcRect/>"
        f'<a:stretch><a:fillRect t="-27000" b="-27000"/></a:stretch>'
        f"</a:blipFill>"
        f"<a:effectLst/>"
        f"</p:bgPr>"
        f"</p:bg>"
    )
    part = slide.part
    r_id = part.get_or_add_image_part(image_path)[1]
    xml = blip_fill.replace("{RID}", r_id)
    from lxml import etree

    bg_el = etree.fromstring(xml)
    cSld = slide._element.find(qn("p:cSld"))
    cSld.insert(0, bg_el)


def add_bar(slide, top_in, height_in, color_hex):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(0), emu(top_in), emu(SLIDE_W_IN), emu(height_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bar_text(slide, top_in, height_in, text, color_hex):
    box = slide.shapes.add_textbox(emu(0), emu(top_in), emu(SLIDE_W_IN), emu(height_in))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(32)
    run.font.color.rgb = RGBColor.from_string(color_hex)
    return box


def add_header_footer(slide, event_name, footer_text, primary_hex, secondary_hex):
    add_bar(slide, 0, 1.05, primary_hex)
    add_bar_text(slide, 0.14, 0.70, event_name, "FFFFFF")
    add_bar(slide, 6.88, 0.67, secondary_hex)
    if footer_text:
        add_bar_text(slide, 6.82, 0.64, footer_text, primary_hex)


def build_title_slide(prs, manifest):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"])
    add_header_footer(
        slide,
        manifest["eventName"],
        manifest["footerText"],
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )

    title_ph = slide.shapes.title
    title_ph.left, title_ph.top = emu(1.67), emu(1.23)
    title_ph.width, title_ph.height = emu(10.00), emu(2.61)
    title_ph.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    for p in title_ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    subtitle_ph = slide.placeholders[1]
    subtitle_ph.left, subtitle_ph.top = emu(1.67), emu(3.94)
    subtitle_ph.width, subtitle_ph.height = emu(10.00), emu(1.81)
    for p in subtitle_ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    return slide


def fit_within(img_path, max_w_in, max_h_in):
    with Image.open(img_path) as im:
        w, h = im.size
    aspect = w / h
    box_aspect = max_w_in / max_h_in
    if aspect > box_aspect:
        draw_w, draw_h = max_w_in, max_w_in / aspect
    else:
        draw_h, draw_w = max_h_in, max_h_in * aspect
    return draw_w, draw_h


def autocrop(image_path):
    """Trims transparent (or near-white) padding around a logo so rasterized
    SVGs -- which Edge renders onto a full-size canvas -- don't end up as a
    tiny mark lost in a mostly-empty grid cell."""
    img = Image.open(image_path)
    if img.mode in ("RGBA", "LA") or (img.mode == "P" and "transparency" in img.info):
        bbox = img.convert("RGBA").split()[-1].getbbox()
    else:
        rgb = img.convert("RGB")
        bg = Image.new("RGB", rgb.size, (255, 255, 255))
        bbox = ImageChops.difference(rgb, bg).getbbox()
    if not bbox or bbox == (0, 0, img.width, img.height):
        return image_path
    cropped_path = str(Path(image_path).with_suffix("")) + "_cropped.png"
    img.crop(bbox).save(cropped_path)
    return cropped_path


def build_sponsor_slide(prs, manifest):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"])
    add_header_footer(
        slide,
        manifest["eventName"],
        "Thank You, Sponsors!",
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )

    sponsors = manifest["sponsors"]
    n = len(sponsors)

    area_left, area_top = 0.35, 1.30
    area_w, area_h = SLIDE_W_IN - 2 * area_left, 6.70 - area_top

    cols = max(1, math.ceil(math.sqrt(n * (area_w / area_h))))
    rows = math.ceil(n / cols)
    while (rows - 1) * cols >= n and rows > 1:
        rows -= 1

    cell_w = area_w / cols
    cell_h = area_h / rows
    pad = 0.12

    for i, sponsor in enumerate(sponsors):
        r, c = divmod(i, cols)
        cell_left = area_left + c * cell_w
        cell_top = area_top + r * cell_h
        logo_path = autocrop(sponsor["logoPath"])
        draw_w, draw_h = fit_within(logo_path, cell_w - 2 * pad, cell_h - 2 * pad)
        pic_left = cell_left + (cell_w - draw_w) / 2
        pic_top = cell_top + (cell_h - draw_h) / 2
        slide.shapes.add_picture(
            logo_path, emu(pic_left), emu(pic_top), emu(draw_w), emu(draw_h)
        )

    return slide


def build_eval_slide(prs, manifest):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(emu(0.92), emu(0.55), emu(11.50), emu(0.87))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Please complete your session evaluations"
    run.font.name = FONT
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])

    # Two-content layout: event logo on the left, an empty QR-code
    # placeholder on the right for the presenter to drop their own
    # Sessionize code into.
    content_top, content_h = 1.60, 4.15
    margin, gap = 0.70, 0.50
    col_w = (SLIDE_W_IN - 2 * margin - gap) / 2
    left_col_x = margin
    right_col_x = margin + col_w + gap

    draw_w, draw_h = fit_within(manifest["logoPath"], col_w, content_h)
    slide.shapes.add_picture(
        manifest["logoPath"],
        emu(left_col_x + (col_w - draw_w) / 2),
        emu(content_top + (content_h - draw_h) / 2),
        emu(draw_w),
        emu(draw_h),
    )

    box_side = min(col_w, content_h)
    box_left = right_col_x + (col_w - box_side) / 2
    box_top = content_top + (content_h - box_side) / 2
    qr_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(box_left), emu(box_top), emu(box_side), emu(box_side)
    )
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    qr_box.line.color.rgb = RGBColor.from_string(manifest["primaryColor"])
    qr_box.line.width = Pt(2.25)
    qr_box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    qr_box.shadow.inherit = False
    qr_tf = qr_box.text_frame
    qr_tf.word_wrap = True
    qr_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    qp = qr_tf.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    r1 = qp.add_run()
    r1.text = "QR CODE"
    r1.font.name = "Calibri"
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])
    qp2 = qr_tf.add_paragraph()
    qp2.alignment = PP_ALIGN.CENTER
    r2 = qp2.add_run()
    r2.text = "Paste your Sessionize evaluation QR code here"
    r2.font.name = "Calibri"
    r2.font.size = Pt(14)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])

    if manifest.get("brugLogoPath"):
        draw_w, draw_h = fit_within(manifest["brugLogoPath"], 1.67, 0.49)
        slide.shapes.add_picture(
            manifest["brugLogoPath"],
            emu((SLIDE_W_IN - draw_w) / 2),
            emu(6.00),
            emu(draw_w),
            emu(draw_h),
        )

    return slide


def build_content_placeholder_slide(prs, manifest):
    """A blank starter slide for presenters to duplicate and build their
    own content on top of."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    box = slide.shapes.add_textbox(emu(0.92), emu(3.10), emu(11.50), emu(1.30))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Your content goes here"
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.italic = True
    run.font.color.rgb = RGBColor.from_string("A6A6A6")

    return slide


def main():
    manifest_path = Path(sys.argv[1])
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = emu(SLIDE_W_IN)
    prs.slide_height = emu(SLIDE_H_IN)

    build_title_slide(prs, manifest)
    build_sponsor_slide(prs, manifest)
    build_content_placeholder_slide(prs, manifest)
    build_eval_slide(prs, manifest)

    out_path = Path(manifest["outputPath"])
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    # Flip the saved package's content type from "presentation" to
    # "template" so it opens in PowerPoint as a .potx (New From Template).
    import zipfile
    import shutil
    import tempfile

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    import os

    os.close(tmp_fd)
    Path(tmp_path).unlink()
    with zipfile.ZipFile(out_path) as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                )
            zout.writestr(item, data)
    shutil.move(tmp_path, out_path)

    print(f"Slide template: {out_path}")


if __name__ == "__main__":
    main()
