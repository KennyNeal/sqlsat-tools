"""
Builds the end-of-day raffle deck (.pptx) as two real PowerPoint sections,
each also saved as a same-named custom show (Slide Show > Custom Slide
Show > Recognition / Raffle):

  Recognition -- self-playing (each slide auto-advances after
    loopAdvanceSeconds): individual slides for the top tiers, grid slides
    for the rest, then a session-evaluation QR slide. Run as a custom show,
    it loops continuously until Esc, so the presenter can run just this
    section as a screensaver before the raffle.
  Raffle -- manually advanced: a "Raffle Time!" intro slide, one hero slide
    per raffle-eligible sponsor (manifest["heroTiers"], shown while that
    sponsor does their drawing), then any raffleDeck.extraHeroSlides (e.g.
    the user group's own drawing, not sourced from sponsors.yaml), then a
    second copy of the evaluation QR slide. Sponsors listed in
    raffleDeck.excludeSponsors (e.g. a sponsor that isn't doing a drawing
    this year) still get their Recognition slide but are skipped here.
    Sponsors listed in raffleDeck.heroLast are pulled out of their normal
    tier position and raffled last, in the order listed -- e.g. to run a
    sponsor's drawing immediately before the extraHeroSlides.

The deck's default "Show slides" range is deliberately left at "All", not
pinned to the Recognition custom show -- pinning it breaks Shift+F5 ("Show
From Current Slide") for every Raffle slide, since they'd fall outside the
allowed range. So the presenter launches each section by name from Slide
Show > Custom Slide Show rather than F5/Shift+F5.

Driven entirely by a JSON manifest (see Generate-RaffleDeck.ps1, which
fetches event/sponsor data and calls this script). Not meant to be run by
hand except for testing -- use the .ps1 wrapper.
"""
import json
import math
import sys
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from slide_helpers import (
    SLIDE_W_IN,
    SLIDE_H_IN,
    FONT,
    emu,
    set_background_watermark,
    add_header_footer,
    add_bar_text,
    fit_within,
    autocrop,
    draw_sponsor_grid,
)


def set_advance_after(slide, seconds):
    """Auto-advances a slide after `seconds` with no click required. Used
    for the self-playing loop section; hero slides skip this call so they
    stay on screen until the presenter clicks forward."""
    from lxml import etree

    transition = etree.SubElement(slide._element, qn("p:transition"))
    transition.set("advTm", str(int(seconds * 1000)))


def build_spotlight_slide(prs, manifest, sponsor, kicker, advance_seconds=None):
    """One sponsor per slide: a kicker line (tier name), a big centered
    logo, and the sponsor name. Used both for the loop deck's individual
    platinum/global slides and for the raffle hero slides."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"])
    add_header_footer(
        slide,
        manifest["eventName"],
        manifest["footerText"],
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )
    add_bar_text(slide, 1.15, 0.55, kicker.upper(), manifest["secondaryColor"], size_pt=24)

    logo_area_top, logo_area_h = 1.85, 3.90
    logo_path = autocrop(sponsor["logoPath"])
    draw_w, draw_h = fit_within(logo_path, 9.0, logo_area_h)
    slide.shapes.add_picture(
        logo_path,
        emu((SLIDE_W_IN - draw_w) / 2),
        emu(logo_area_top + (logo_area_h - draw_h) / 2),
        emu(draw_w),
        emu(draw_h),
    )

    add_bar_text(slide, 5.85, 0.85, sponsor["name"], manifest["primaryColor"], size_pt=36)

    if advance_seconds:
        set_advance_after(slide, advance_seconds)

    return slide


def build_grid_slide(prs, manifest, sponsors, heading, advance_seconds=None):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"])
    add_header_footer(
        slide,
        manifest["eventName"],
        manifest["footerText"],
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )
    add_bar_text(slide, 1.10, 0.55, heading, manifest["primaryColor"], size_pt=24)

    area_left, area_top = 0.35, 1.80
    area_w, area_h = SLIDE_W_IN - 2 * area_left, 6.70 - area_top
    draw_sponsor_grid(slide, sponsors, area_left, area_top, area_w, area_h)

    if advance_seconds:
        set_advance_after(slide, advance_seconds)

    return slide


def build_eval_slide(prs, manifest, advance_seconds=None):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"])
    add_header_footer(
        slide,
        manifest["eventName"],
        manifest["footerText"],
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )

    title_box = slide.shapes.add_textbox(emu(0.92), emu(1.30), emu(11.50), emu(0.87))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Please Evaluate Today's Sessions"
    run.font.name = FONT
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])

    qr_side = 3.4
    qr_left = (SLIDE_W_IN - qr_side) / 2
    qr_top = 2.35
    slide.shapes.add_picture(
        manifest["qrPath"], emu(qr_left), emu(qr_top), emu(qr_side), emu(qr_side)
    )

    cap_box = slide.shapes.add_textbox(
        emu(0.92), emu(qr_top + qr_side + 0.15), emu(11.50), emu(0.5)
    )
    ctf = cap_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = manifest["evalUrl"]
    cr.font.name = "Calibri"
    cr.font.size = Pt(16)
    cr.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])

    if advance_seconds:
        set_advance_after(slide, advance_seconds)

    return slide


def build_raffle_intro_slide(prs, manifest, text, subtext=None):
    """The first slide of the Raffle section: a big banner the presenter
    lands on right after escaping the looping Recognition custom show."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"], alpha_pct=8)
    add_header_footer(
        slide,
        manifest["eventName"],
        manifest["footerText"],
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )

    box = slide.shapes.add_textbox(emu(0.92), emu(2.80), emu(11.50), emu(1.30))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])

    if subtext:
        box2 = slide.shapes.add_textbox(emu(0.92), emu(4.05), emu(11.50), emu(0.6))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtext
        r2.font.name = "Calibri"
        r2.font.size = Pt(16)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor.from_string("808080")

    return slide


def balanced_chunks(items, max_size):
    """Splits into as-even-as-possible groups no larger than max_size, so a
    tier just over the limit (e.g. 9 sponsors, max 8) doesn't leave a
    near-empty trailing slide with a single lonely logo."""
    n = len(items)
    if n == 0:
        return []
    num_chunks = math.ceil(n / max_size)
    base, extra = divmod(n, num_chunks)
    chunks, idx = [], 0
    for i in range(num_chunks):
        size = base + (1 if i < extra else 0)
        chunks.append(items[idx : idx + size])
        idx += size
    return chunks


SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


def _custom_show_list_xml(shows):
    """shows: ordered [(name, show_id, [r_id, ...]), ...]"""
    show_els = []
    for name, show_id, r_ids in shows:
        slides_xml = "".join(f'<p:sld r:id="{rid}"/>' for rid in r_ids)
        show_els.append(
            f'<p:custShow name="{name}" id="{show_id}">'
            f"<p:sldLst>{slides_xml}</p:sldLst></p:custShow>"
        )
    return f'<p:custShowLst>{"".join(show_els)}</p:custShowLst>'


def _section_list_xml(sections):
    """sections: ordered [(name, [slide_id, ...]), ...]"""
    section_els = []
    for name, slide_ids in sections:
        guid = "{" + str(uuid.uuid4()).upper() + "}"
        ids_xml = "".join(f'<p14:sldId id="{sid}"/>' for sid in slide_ids)
        section_els.append(
            f'<p14:section name="{name}" id="{guid}">'
            f"<p14:sldIdLst>{ids_xml}</p14:sldIdLst></p14:section>"
        )
    return (
        f'<p:extLst><p:ext uri="{SECTION_EXT_URI}">'
        f'<p14:sectionLst xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">'
        f'{"".join(section_els)}</p14:sectionLst></p:ext></p:extLst>'
    )


def set_show_setup(pptx_path, shows, sections):
    """Adds real PowerPoint sections and two named custom shows
    ("Recognition", "Raffle"), and turns on loop-until-Esc. python-pptx has
    no API for any of this: sections are a PowerPoint-2010 extension
    (p14:sectionLst in presentation.xml's extLst), and show properties live
    in the separate presProps.xml part. So it's a direct post-save zip edit
    -- same technique generate_slide_template.py uses to flip a .pptx into
    a .potx -- against the exact presentation.xml layout python-pptx's
    default template produces (a <p:defaultTextStyle> immediately after
    <p:notesSz>, and no pre-existing <p:extLst>).

    The default "Show slides" range is deliberately left at "All" (no
    <p:custShow> in showPr) rather than pinned to Recognition: pinning it
    breaks "Show From Current Slide" (Shift+F5) for any slide outside that
    custom show, which is exactly how a presenter starts the Raffle
    section. Instead, both sections are launched by name from Slide Show >
    Custom Slide Show; loop applies to whichever one is currently running."""
    import zipfile
    import shutil
    import tempfile
    import os

    custom_show_xml = _custom_show_list_xml(shows)
    section_list_xml = _section_list_xml(sections)
    show_pr_xml = '<p:showPr loop="1"><p:present/></p:showPr>'

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(tmp_fd)
    Path(tmp_path).unlink()
    with zipfile.ZipFile(pptx_path) as zin, zipfile.ZipFile(
        tmp_path, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/presentation.xml":
                marker = b"<p:defaultTextStyle>"
                if marker not in data:
                    raise RuntimeError(
                        "presentation.xml layout changed; can't insert p:custShowLst"
                    )
                data = data.replace(marker, custom_show_xml.encode("utf-8") + marker, 1)
                data = data.replace(
                    b"</p:presentation>",
                    section_list_xml.encode("utf-8") + b"</p:presentation>",
                    1,
                )
            elif item.filename == "ppt/presProps.xml":
                if b"<p:extLst>" in data:
                    data = data.replace(b"<p:extLst>", show_pr_xml.encode("utf-8") + b"<p:extLst>", 1)
                else:
                    data = data.replace(
                        b"</p:presentationPr>",
                        show_pr_xml.encode("utf-8") + b"</p:presentationPr>",
                        1,
                    )
            zout.writestr(item, data)
    shutil.move(tmp_path, pptx_path)


def main():
    manifest_path = Path(sys.argv[1])
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = emu(SLIDE_W_IN)
    prs.slide_height = emu(SLIDE_H_IN)

    by_tier = {g["tier"]: g for g in manifest["groups"]}
    advance_seconds = manifest["loopAdvanceSeconds"]

    # ── Recognition section (self-playing, advances on a timer) ─────────
    recognition_slides = []

    for tier in manifest["individualTiers"]:
        group = by_tier.get(tier)
        if not group:
            continue
        for sponsor in group["sponsors"]:
            recognition_slides.append(
                build_spotlight_slide(prs, manifest, sponsor, group["title"], advance_seconds)
            )

    for tier_group in manifest["gridGroups"]:
        sponsors, titles = [], []
        for tier in tier_group:
            group = by_tier.get(tier)
            if not group:
                continue
            sponsors.extend(group["sponsors"])
            titles.append(group["title"])
        if not sponsors:
            continue
        heading = " & ".join(titles)
        for batch in balanced_chunks(sponsors, manifest["maxPerGridSlide"]):
            recognition_slides.append(
                build_grid_slide(prs, manifest, batch, heading, advance_seconds)
            )

    recognition_slides.append(build_eval_slide(prs, manifest, advance_seconds))

    # ── Raffle section (manual advance) ──────────────────────────────────
    raffle_slides = [
        build_raffle_intro_slide(
            prs,
            manifest,
            "\U0001F389 Raffle Time!",
            "Presenter: advance manually from here through the prize drawing",
        )
    ]

    excluded = set(manifest.get("excludeSponsors", []))
    hero_last = manifest.get("heroLast", [])
    deferred = {}  # sponsor name -> (sponsor, tier title), raffled last, in heroLast order
    for tier in manifest["heroTiers"]:
        group = by_tier.get(tier)
        if not group:
            continue
        for sponsor in group["sponsors"]:
            if sponsor["name"] in excluded:
                print(f"  Skipping raffle hero slide (excluded): {sponsor['name']}")
                continue
            if sponsor["name"] in hero_last:
                deferred[sponsor["name"]] = (sponsor, group["title"])
                continue
            raffle_slides.append(build_spotlight_slide(prs, manifest, sponsor, group["title"]))

    for name in hero_last:
        entry = deferred.get(name)
        if not entry:
            print(f"  Warning: heroLast sponsor '{name}' not found among heroTiers")
            continue
        sponsor, title = entry
        raffle_slides.append(build_spotlight_slide(prs, manifest, sponsor, title))

    for extra in manifest.get("extraHeroSlides", []):
        sponsor = {"name": extra["name"], "logoPath": extra["logoPath"]}
        raffle_slides.append(build_spotlight_slide(prs, manifest, sponsor, extra["kicker"]))

    raffle_slides.append(build_eval_slide(prs, manifest))

    # ── Sections + named custom shows ("Recognition", "Raffle") ─────────
    id_to_rid = {
        int(el.get("id")): el.get(qn("r:id")) for el in prs.slides._sldIdLst
    }
    shows = [
        ("Recognition", 0, [id_to_rid[slide.slide_id] for slide in recognition_slides]),
        ("Raffle", 1, [id_to_rid[slide.slide_id] for slide in raffle_slides]),
    ]
    sections = [
        ("Recognition", [slide.slide_id for slide in recognition_slides]),
        ("Raffle", [slide.slide_id for slide in raffle_slides]),
    ]

    out_path = Path(manifest["outputPath"])
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    set_show_setup(out_path, shows, sections)

    print(f"Raffle deck: {out_path}")


if __name__ == "__main__":
    main()
