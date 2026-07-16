"""
Builds the presenter slide-deck template (.potx): title slide, sponsor
"thank you" slide, and session-evaluation slide.

Driven entirely by a JSON manifest (see Generate-SlideTemplate.ps1, which
fetches event/sponsor data and calls this script). Not meant to be run by
hand except for testing -- use the .ps1 wrapper.
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

from slide_helpers import (
    SLIDE_W_IN,
    SLIDE_H_IN,
    FONT,
    emu,
    set_background_watermark,
    add_header_footer,
    fit_within,
    draw_sponsor_grid,
)


def build_title_slide(prs, manifest):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"])
    add_header_footer(
        slide,
        manifest["eventName"],
        manifest["footerText"],
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )

    title_ph = slide.shapes.title
    title_ph.left, title_ph.top = emu(1.67), emu(1.23)
    title_ph.width, title_ph.height = emu(10.00), emu(2.61)
    title_ph.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    for p in title_ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    subtitle_ph = slide.placeholders[1]
    subtitle_ph.left, subtitle_ph.top = emu(1.67), emu(3.94)
    subtitle_ph.width, subtitle_ph.height = emu(10.00), emu(1.81)
    for p in subtitle_ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    return slide


def build_sponsor_slide(prs, manifest):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    set_background_watermark(slide, manifest["logoPath"])
    add_header_footer(
        slide,
        manifest["eventName"],
        "Thank You, Sponsors!",
        manifest["primaryColor"],
        manifest["secondaryColor"],
    )

    area_left, area_top = 0.35, 1.30
    area_w, area_h = SLIDE_W_IN - 2 * area_left, 6.70 - area_top
    draw_sponsor_grid(slide, manifest["sponsors"], area_left, area_top, area_w, area_h)

    return slide


def build_eval_slide(prs, manifest):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(emu(0.92), emu(0.55), emu(11.50), emu(0.87))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Please complete your session evaluations"
    run.font.name = FONT
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])

    # Two-content layout: event logo on the left, an empty QR-code
    # placeholder on the right for the presenter to drop their own
    # Sessionize code into.
    content_top, content_h = 1.60, 4.15
    margin, gap = 0.70, 0.50
    col_w = (SLIDE_W_IN - 2 * margin - gap) / 2
    left_col_x = margin
    right_col_x = margin + col_w + gap

    draw_w, draw_h = fit_within(manifest["logoPath"], col_w, content_h)
    slide.shapes.add_picture(
        manifest["logoPath"],
        emu(left_col_x + (col_w - draw_w) / 2),
        emu(content_top + (content_h - draw_h) / 2),
        emu(draw_w),
        emu(draw_h),
    )

    box_side = min(col_w, content_h)
    box_left = right_col_x + (col_w - box_side) / 2
    box_top = content_top + (content_h - box_side) / 2
    qr_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(box_left), emu(box_top), emu(box_side), emu(box_side)
    )
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    qr_box.line.color.rgb = RGBColor.from_string(manifest["primaryColor"])
    qr_box.line.width = Pt(2.25)
    qr_box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    qr_box.shadow.inherit = False
    qr_tf = qr_box.text_frame
    qr_tf.word_wrap = True
    qr_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    qp = qr_tf.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    r1 = qp.add_run()
    r1.text = "QR CODE"
    r1.font.name = "Calibri"
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])
    qp2 = qr_tf.add_paragraph()
    qp2.alignment = PP_ALIGN.CENTER
    r2 = qp2.add_run()
    r2.text = "Paste your Sessionize evaluation QR code here"
    r2.font.name = "Calibri"
    r2.font.size = Pt(14)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor.from_string(manifest["primaryColor"])

    if manifest.get("brugLogoPath"):
        draw_w, draw_h = fit_within(manifest["brugLogoPath"], 1.67, 0.49)
        slide.shapes.add_picture(
            manifest["brugLogoPath"],
            emu((SLIDE_W_IN - draw_w) / 2),
            emu(6.00),
            emu(draw_w),
            emu(draw_h),
        )

    return slide


def build_content_placeholder_slide(prs, manifest):
    """A blank starter slide for presenters to duplicate and build their
    own content on top of."""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    box = slide.shapes.add_textbox(emu(0.92), emu(3.10), emu(11.50), emu(1.30))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Your content goes here"
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.italic = True
    run.font.color.rgb = RGBColor.from_string("A6A6A6")

    return slide


def main():
    manifest_path = Path(sys.argv[1])
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = emu(SLIDE_W_IN)
    prs.slide_height = emu(SLIDE_H_IN)

    build_title_slide(prs, manifest)
    build_sponsor_slide(prs, manifest)
    build_content_placeholder_slide(prs, manifest)
    build_eval_slide(prs, manifest)

    out_path = Path(manifest["outputPath"])
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    # Flip the saved package's content type from "presentation" to
    # "template" so it opens in PowerPoint as a .potx (New From Template).
    import zipfile
    import shutil
    import tempfile

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    import os

    os.close(tmp_fd)
    Path(tmp_path).unlink()
    with zipfile.ZipFile(out_path) as zin, zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                )
            zout.writestr(item, data)
    shutil.move(tmp_path, out_path)

    print(f"Slide template: {out_path}")


if __name__ == "__main__":
    main()
