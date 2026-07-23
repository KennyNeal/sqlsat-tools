"""
Shared PowerPoint-building helpers used by generate_slide_template.py and
generate_raffle_deck.py: EMU/inch conversion, branded header/footer bars, the
faded event-logo background watermark, image fitting/autocropping, and the
sponsor-logo grid layout.
"""
import math
from pathlib import Path

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageChops

EMU_PER_IN = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

FONT = "Berlin Sans FB"


def emu(inches):
    return Emu(int(round(inches * EMU_PER_IN)))


def set_background_watermark(slide, image_path, alpha_pct=12):
    """Sets a slide's background to a picture, faded to alpha_pct opacity,
    stretched/cropped the way the previous year's template did (slightly
    cropped top and bottom so the logo reads large and centered)."""
    blip_fill = (
        f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f"<p:bgPr>"
        f'<a:blipFill dpi="0" rotWithShape="1">'
        f'<a:blip r:embed="{{RID}}"><a:alphaModFix amt="{alpha_pct * 1000}"/></a:blip>'
        f"<a:srcRect/>"
        f'<a:stretch><a:fillRect t="-27000" b="-27000"/></a:stretch>'
        f"</a:blipFill>"
        f"<a:effectLst/>"
        f"</p:bgPr>"
        f"</p:bg>"
    )
    part = slide.part
    r_id = part.get_or_add_image_part(image_path)[1]
    xml = blip_fill.replace("{RID}", r_id)
    from lxml import etree

    bg_el = etree.fromstring(xml)
    cSld = slide._element.find(qn("p:cSld"))
    cSld.insert(0, bg_el)


def add_bar(slide, top_in, height_in, color_hex):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(0), emu(top_in), emu(SLIDE_W_IN), emu(height_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bar_text(slide, top_in, height_in, text, color_hex, size_pt=32):
    box = slide.shapes.add_textbox(emu(0), emu(top_in), emu(SLIDE_W_IN), emu(height_in))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size_pt)
    run.font.color.rgb = RGBColor.from_string(color_hex)
    return box


def add_header_footer(slide, event_name, footer_text, primary_hex, secondary_hex):
    add_bar(slide, 0, 1.05, primary_hex)
    add_bar_text(slide, 0.14, 0.70, event_name, "FFFFFF")
    add_bar(slide, 6.88, 0.67, secondary_hex)
    if footer_text:
        add_bar_text(slide, 6.82, 0.64, footer_text, primary_hex)


def fit_within(img_path, max_w_in, max_h_in):
    with Image.open(img_path) as im:
        w, h = im.size
    aspect = w / h
    box_aspect = max_w_in / max_h_in
    if aspect > box_aspect:
        draw_w, draw_h = max_w_in, max_w_in / aspect
    else:
        draw_h, draw_w = max_h_in, max_h_in * aspect
    return draw_w, draw_h


def autocrop(image_path):
    """Trims transparent (or near-white) padding around a logo so rasterized
    SVGs -- which Edge renders onto a full-size canvas -- don't end up as a
    tiny mark lost in a mostly-empty grid cell."""
    img = Image.open(image_path)
    if img.mode in ("RGBA", "LA") or (img.mode == "P" and "transparency" in img.info):
        bbox = img.convert("RGBA").split()[-1].getbbox()
    else:
        rgb = img.convert("RGB")
        bg = Image.new("RGB", rgb.size, (255, 255, 255))
        bbox = ImageChops.difference(rgb, bg).getbbox()
    if not bbox or bbox == (0, 0, img.width, img.height):
        return image_path
    cropped_path = str(Path(image_path).with_suffix("")) + "_cropped.png"
    img.crop(bbox).save(cropped_path)
    return cropped_path


def _balanced_row_sizes(n, rows):
    """Splits n items across `rows` rows as evenly as possible (e.g. 5 items
    over 2 rows -> [3, 2]), so no row ends up sparse next to full ones."""
    base, extra = divmod(n, rows)
    return [base + 1 if r < extra else base for r in range(rows)]


def draw_sponsor_grid(slide, sponsors, area_left, area_top, area_w, area_h):
    """Lays sponsor logos out in a roughly-square grid sized to fill the
    given area, autocropping and fitting each logo within its cell. Each
    row's item count is balanced (not a fixed column count sliced
    row-major), and every row's cells stretch to the full area width, so a
    trailing row with fewer sponsors doesn't leave a dead gap -- it gets
    fewer, larger cells instead."""
    n = len(sponsors)
    approx_cols = max(1, round(math.sqrt(n * (area_w / area_h))))
    rows = max(1, math.ceil(n / approx_cols))
    row_sizes = _balanced_row_sizes(n, rows)

    cell_h = area_h / rows
    pad = 0.12

    idx = 0
    for r, row_n in enumerate(row_sizes):
        cell_w = area_w / row_n
        row_top = area_top + r * cell_h
        for c in range(row_n):
            sponsor = sponsors[idx]
            idx += 1
            cell_left = area_left + c * cell_w
            logo_path = autocrop(sponsor["logoPath"])
            draw_w, draw_h = fit_within(logo_path, cell_w - 2 * pad, cell_h - 2 * pad)
            pic_left = cell_left + (cell_w - draw_w) / 2
            pic_top = row_top + (cell_h - draw_h) / 2
            slide.shapes.add_picture(
                logo_path, emu(pic_left), emu(pic_top), emu(draw_w), emu(draw_h)
            )
